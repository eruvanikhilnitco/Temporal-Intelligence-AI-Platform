import logging
import re
import json
import os
from pathlib import Path
from typing import List, Dict, Optional

import spacy

logger = logging.getLogger(__name__)

# Configure Tika to use the pre-downloaded JAR and disable verbose output
os.environ.setdefault("TIKA_SERVER_JAR", "/tmp/tika-server.jar")
os.environ.setdefault("TIKA_LOG_PATH", "/tmp")
os.environ.setdefault("TIKA_STARTUP_MAX_RETRY", "1")  # fail fast

# Formats handled natively — never touches Tika
_NATIVE_EXTS = {".txt", ".md", ".html", ".htm", ".xml", ".csv", ".json", ".log"}

# Formats handled with pure-Python libs — no JVM needed
_PYTHON_EXTS = {".pdf", ".docx", ".doc", ".pptx", ".xlsx", ".xls"}


def _extract_text_native(path: Path) -> str:
    """Read plain text files directly."""
    try:
        return path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        logger.warning("Native read failed for %s: %s", path, e)
        return ""


def _extract_text_pdf(path: Path) -> str:
    """Extract text from PDF using pypdf (pure Python, no Java)."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        parts = []
        for page in reader.pages:
            try:
                t = page.extract_text()
                if t:
                    parts.append(t)
            except Exception:
                pass
        return "\n".join(parts)
    except Exception as e:
        logger.warning("pypdf failed for %s: %s — trying Tika", path, e)
        return _extract_text_tika(path)


def _extract_text_docx(path: Path) -> str:
    """Extract text from DOCX using python-docx."""
    try:
        import docx as _docx
        doc = _docx.Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        logger.warning("python-docx failed for %s: %s — trying Tika", path, e)
        return _extract_text_tika(path)


def _extract_text_pptx(path: Path) -> str:
    """Extract text from PPTX using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
        return "\n".join(parts)
    except Exception as e:
        logger.warning("python-pptx failed for %s: %s — trying Tika", path, e)
        return _extract_text_tika(path)


def _extract_text_xlsx(path: Path) -> str:
    """Extract text from XLSX/XLS using openpyxl."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as e:
        logger.warning("openpyxl failed for %s: %s — trying Tika", path, e)
        return _extract_text_tika(path)


def _extract_text_csv(path: Path) -> str:
    """Parse CSV natively."""
    try:
        import csv
        rows = []
        with open(path, newline="", encoding="utf-8", errors="replace") as f:
            for i, row in enumerate(csv.reader(f)):
                if i >= 5000:
                    break
                rows.append(" | ".join(row))
        return "\n".join(rows)
    except Exception as e:
        logger.warning("CSV parse failed for %s: %s", path, e)
        return _extract_text_native(path)


def _extract_text_json(path: Path) -> str:
    """Parse JSON into readable text."""
    try:
        with open(path, encoding="utf-8", errors="replace") as f:
            data = json.load(f)
        if isinstance(data, dict):
            return "\n".join(f"{k}: {v}" for k, v in data.items())
        if isinstance(data, list):
            return "\n".join(str(item) for item in data[:2000])
        return str(data)
    except Exception:
        return _extract_text_native(path)


def _extract_text_tika(path: Path) -> str:
    """
    Last-resort fallback: Apache Tika via Java.
    Only called for formats that have no pure-Python alternative.
    """
    try:
        from tika import parser as tika_parser
        parsed = tika_parser.from_file(str(path))
        text = (parsed or {}).get("content", "") or ""
        if not text.strip():
            raise ValueError("Tika returned empty content")
        return text
    except Exception as e:
        logger.error("Tika failed for %s: %s — falling back to raw text read", path, e)
        try:
            return path.read_text(encoding="utf-8", errors="replace")
        except Exception:
            return f"Document: {path.name}"


class DocumentParser:
    def __init__(self, spacy_model: str = "en_core_web_sm"):
        try:
            self.nlp = spacy.load(spacy_model)
        except OSError:
            logger.warning("spaCy model %s not found — entity extraction disabled", spacy_model)
            self.nlp = None

    def _extract_text(self, path: Path) -> str:
        """Route to the right extractor based on file extension."""
        ext = path.suffix.lower()

        if ext in _NATIVE_EXTS:
            if ext == ".csv":
                return _extract_text_csv(path)
            if ext == ".json":
                return _extract_text_json(path)
            return _extract_text_native(path)

        if ext == ".pdf":
            return _extract_text_pdf(path)
        if ext in (".docx", ".doc"):
            return _extract_text_docx(path)
        if ext in (".pptx", ".ppt"):
            return _extract_text_pptx(path)
        if ext in (".xlsx", ".xls"):
            return _extract_text_xlsx(path)

        # Unknown binary format — try Tika
        return _extract_text_tika(path)

    def parse(self, file_path: str, original_filename: Optional[str] = None) -> dict:
        """
        Parse document without blocking on Tika for text-based formats.
        Returns: {text, entities, detected_document_type, detected_lifecycle_id}
        """
        path = Path(file_path)
        filename = original_filename or path.name

        text = self._extract_text(path)
        if not text or not text.strip():
            text = f"Document: {filename}"

        # Document type detection
        detected_doc_type = self._detect_document_type(text, filename)
        detected_lifecycle_id = self._extract_lifecycle_id(text)

        # Entity extraction with spaCy
        entities: List[str] = []
        if self.nlp and text:
            try:
                doc = self.nlp(text[:100_000])
                seen: set = set()
                for ent in doc.ents:
                    if ent.label_ in ("ORG", "PERSON", "MONEY", "DATE", "PRODUCT"):
                        if ent.text not in seen:
                            seen.add(ent.text)
                            entities.append(ent.text)
                            if len(entities) >= 50:
                                break
            except Exception as e:
                logger.error("spaCy entity extraction failed: %s", e)

        return {
            "text": text.strip(),
            "entities": entities,
            "detected_document_type": detected_doc_type,
            "detected_lifecycle_id": detected_lifecycle_id,
        }

    # ── Document type detection ─────────────────────────────────────────────

    def _detect_document_type(self, text: str, filename: str) -> str:
        text_lower = text.lower()
        filename_lower = filename.lower()

        patterns = {
            "Patient Record": [r"patient\s*record\b", r"medical\s*record\b", r"vital\s*signs", r"diagnosis\s*:"],
            "Medical Report": [r"medical\s*report\b", r"clinical\s*report", r"lab\s*result"],
            "Prescription": [r"prescription\b", r"medication\s*order\b", r"rx\b"],
            "Resume": [r"resume\b", r"\bcv\b", r"curriculum\s*vitae", r"professional\s*summary"],
            "Application": [r"application\s*form\b", r"job\s*application\b", r"candidate\s*:"],
            "Offer Letter": [r"offer\s*letter\b", r"job\s*offer\b", r"employment\s*offer\b"],
            "Financial Statement": [r"financial\s*statement\b", r"financial\s*report"],
            "Compliance Report": [r"compliance\s*report\b", r"regulatory\s*compliance"],
            "Expense Report": [r"expense\s*report\b", r"expense\s*statement"],
            "Invoice": [r"\binvoice\b", r"invoice\s*number", r"bill\s*to\b"],
            "Purchase Order": [r"purchase\s*order\b", r"\bpo\s*#\b"],
            "Contract": [r"contract\b", r"agreement\b", r"terms\s*and\s*conditions"],
            "Report": [r"\breport\b", r"\bsummary\b", r"\banalysis\b"],
            "Certificate": [r"certificate\b", r"certification\b"],
        }

        for doc_type, pats in patterns.items():
            for pat in pats:
                if re.search(pat, filename_lower, re.IGNORECASE):
                    return doc_type
        for doc_type, pats in patterns.items():
            for pat in pats:
                if re.search(pat, text_lower, re.IGNORECASE):
                    return doc_type

        ext_map = {".pdf": "PDF Document", ".docx": "Word Document", ".doc": "Word Document",
                   ".txt": "Text Document", ".csv": "CSV Data", ".xlsx": "Excel Spreadsheet",
                   ".xls": "Excel Spreadsheet", ".json": "JSON Data", ".xml": "XML Document",
                   ".pptx": "Presentation", ".md": "Markdown Document"}
        return ext_map.get(Path(filename).suffix.lower(), "Document")

    def _extract_lifecycle_id(self, text: str) -> Optional[str]:
        patterns = [
            r"lifecycle[_\s-]?id[:\s]+([A-Za-z0-9_-]+)",
            r"lifecycle[:\s]+([A-Za-z0-9_-]+)",
            r"\b(lifecycle[_\s-]?[0-9]{3,})\b",
            r"invoice[_\s-]?#?[:\s]+([A-Za-z0-9_-]+)",
            r"contract[_\s-]?#?[:\s]+([A-Za-z0-9_-]+)",
        ]
        sample = "\n".join(text.split("\n")[:50])
        for pat in patterns:
            m = re.search(pat, sample, re.IGNORECASE | re.MULTILINE)
            if m:
                lid = re.sub(r"[^\w_-]", "", m.group(1).strip())
                if len(lid) >= 3:
                    return lid.lower()
        return None
